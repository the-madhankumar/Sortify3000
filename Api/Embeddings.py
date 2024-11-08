import os
import pdfplumber
from pptx import Presentation
import pandas as pd
from langchain.embeddings import HuggingFaceBgeEmbeddings
from transformers import AutoTokenizer
from VectorData import store_embedding

embeddings = HuggingFaceBgeEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
tokenizer = AutoTokenizer.from_pretrained("sentence-transformers/all-MiniLM-L6-v2")

def read_document(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    text = ""

    if ext == ".pdf":
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
    elif ext == ".pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + ' '
    elif ext in ['.xls', '.xlsx']:
        df = pd.read_excel(file_path)
        text = df.to_string(index=False)
    else:
        with open(file_path, 'r') as file:
            text = file.read() 
    
    return text

def chunk_text(text, max_length=512):
    tokens = tokenizer.encode(text, truncation=False)
    chunks = [tokens[i:i+max_length] for i in range(0, len(tokens), max_length)]
    return [tokenizer.decode(chunk) for chunk in chunks]

file_path = r"C:\Users\madha\Downloads\qbmin.pdf" 

document_text = read_document(file_path)
chunked_texts = chunk_text(document_text, max_length=512)

for idx, chunk in enumerate(chunked_texts):
    text_embedding = embeddings.embed_documents([chunk])[0]
    doc_id = f"doc_{os.path.basename(file_path)}_chunk_{idx}"
    metadata = {"file_path": file_path, "chunk_index": idx}
    
    store_embedding(text_embedding, doc_id, metadata)

print("All embeddings stored successfully.")
